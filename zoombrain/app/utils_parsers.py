# Attachment -> text helpers (PDF/PPTX/TXT)

from typing import Optional
import os
from pathlib import Path


def extract_text_from_pdf(file_path: str) -> str:
    """
    Extract text from PDF file
    
    Args:
        file_path: Path to PDF file
        
    Returns:
        Extracted text as string
    """
    try:
        import PyPDF2
        
        text = ""
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
        
        return text.strip()
    except Exception as e:
        raise Exception(f"Error extracting text from PDF: {str(e)}")


def extract_text_from_pptx(file_path: str) -> str:
    """
    Extract text from PowerPoint file
    
    Args:
        file_path: Path to PPTX file
        
    Returns:
        Extracted text as string
    """
    try:
        from pptx import Presentation
        
        text = ""
        prs = Presentation(file_path)
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        
        return text.strip()
    except Exception as e:
        raise Exception(f"Error extracting text from PPTX: {str(e)}")


def extract_text_from_txt(file_path: str, encoding: str = 'utf-8') -> str:
    """
    Extract text from text file
    
    Args:
        file_path: Path to text file
        encoding: Text encoding (default: utf-8)
        
    Returns:
        File content as string
    """
    try:
        with open(file_path, 'r', encoding=encoding) as file:
            return file.read().strip()
    except Exception as e:
        raise Exception(f"Error reading text file: {str(e)}")


def extract_text_from_file(file_path: str) -> str:
    """
    Extract text from file based on extension
    
    Args:
        file_path: Path to file
        
    Returns:
        Extracted text as string
    """
    extension = Path(file_path).suffix.lower()
    
    extractors = {
        '.pdf': extract_text_from_pdf,
        '.pptx': extract_text_from_pptx,
        '.txt': extract_text_from_txt,
        '.md': extract_text_from_txt,
    }
    
    if extension not in extractors:
        raise ValueError(f"Unsupported file type: {extension}")
    
    return extractors[extension](file_path)


def chunk_text(text: str, chunk_size: int = 500, overlap: int = 50) -> list[str]:
    """
    Split text into overlapping chunks
    
    Args:
        text: Input text
        chunk_size: Size of each chunk in characters
        overlap: Overlap between chunks
        
    Returns:
        List of text chunks
    """
    chunks = []
    start = 0
    
    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end]
        chunks.append(chunk)
        start += chunk_size - overlap
    
    return chunks
